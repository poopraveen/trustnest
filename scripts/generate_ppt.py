from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ─── Brand Colours ──────────────────────────────────────────────────────────
GREEN      = RGBColor(0x16, 0x65, 0x34)   # dark green
GREEN_MID  = RGBColor(0x22, 0xC5, 0x5E)   # emerald
TEAL       = RGBColor(0x0D, 0x94, 0x88)   # teal accent
DARK       = RGBColor(0x0F, 0x17, 0x2A)   # near-black
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF0, 0xFD, 0xF4)   # light green bg
AMBER      = RGBColor(0xD9, 0x77, 0x06)
SLATE      = RGBColor(0x47, 0x55, 0x69)
SILVER     = RGBColor(0xE2, 0xE8, 0xF0)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]   # completely blank

# ─── Helpers ────────────────────────────────────────────────────────────────
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic
    return tb

def card(slide, l, t, w, h, fill=WHITE, border=SILVER):
    r = rect(slide, l, t, w, h, fill, border, 0.5)
    # rounded corners via XML
    sp = r._element
    spPr = sp.find(nsmap.qn('p:spPr'))
    prstGeom = spPr.find(nsmap.qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'roundRect')
        avLst = prstGeom.find(nsmap.qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, nsmap.qn('a:avLst'))
        avLst.clear()
        gd = etree.SubElement(avLst, nsmap.qn('a:gd'))
        gd.set('name','adj'); gd.set('fmla','val 30000')
    return r

def bullet_list(slide, items, l, t, w, size, color, spacing=0.38, icon="✦"):
    for i, item in enumerate(items):
        txt(slide, f"{icon}  {item}", l, t + i*spacing, w, spacing, size, color)

def gradient_rect(slide, l, t, w, h, c1, c2):
    """Simulate gradient with two overlapping rects (left→right)."""
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = c1
    shape.line.fill.background()
    # overlay a semi-transparent dark on right half
    return shape

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)

# Green accent bar left
rect(s, 0, 0, 0.18, 7.5, GREEN)
# Teal accent bottom strip
rect(s, 0.18, 7.0, 13.15, 0.5, TEAL)

# Large circle decoration (top right)
circle = s.shapes.add_shape(9, Inches(9.5), Inches(-1.2), Inches(4.5), Inches(4.5))
circle.fill.solid(); circle.fill.fore_color.rgb = GREEN
circle.line.fill.background()
# opacity via XML
sp = circle._element
spPr = sp.find(nsmap.qn('p:spPr'))
solidFill = spPr.find('.//' + nsmap.qn('a:solidFill'))
if solidFill is not None:
    srgb = solidFill.find(nsmap.qn('a:srgbClr'))
    if srgb is not None:
        lum = etree.SubElement(srgb, nsmap.qn('a:alpha'))
        lum.set('val','15000')

# Small circle
circle2 = s.shapes.add_shape(9, Inches(10.8), Inches(4.5), Inches(2.2), Inches(2.2))
circle2.fill.solid(); circle2.fill.fore_color.rgb = TEAL
circle2.line.fill.background()
sp2 = circle2._element
spPr2 = sp2.find(nsmap.qn('p:spPr'))
sf2 = spPr2.find('.//' + nsmap.qn('a:solidFill'))
if sf2:
    srgb2 = sf2.find(nsmap.qn('a:srgbClr'))
    if srgb2 is not None:
        a2 = etree.SubElement(srgb2, nsmap.qn('a:alpha')); a2.set('val','12000')

txt(s, "🌿", 0.5, 0.6, 1.5, 1.0, 36, GREEN_MID, bold=True, align=PP_ALIGN.CENTER)
txt(s, "TrustNest", 0.4, 1.55, 8, 1.2, 68, WHITE, bold=True)
txt(s, "HBL Club Management Platform", 0.4, 2.65, 9, 0.7, 28, GREEN_MID, bold=True)
txt(s, "The All-In-One Digital Solution for Herbalife Nutrition Clubs", 0.4, 3.3, 9.5, 0.5, 16, SILVER, italic=True)

# Stats strip
for i,(val,lbl) in enumerate([("500+","Members Managed"),("3x","Retention Lift"),("100%","Paperless"),("₹0","Setup Cost")]):
    xoff = 0.4 + i*3.0
    rect(s, xoff, 5.1, 2.6, 1.3, RGBColor(0x1a,0x2a,0x3a), TEAL, 0.5)
    txt(s, val, xoff+0.1, 5.18, 2.4, 0.55, 26, GREEN_MID, bold=True, align=PP_ALIGN.CENTER)
    txt(s, lbl, xoff+0.1, 5.68, 2.4, 0.35, 10, SILVER, align=PP_ALIGN.CENTER)

txt(s, "Confidential Pitch Deck  •  2026", 0.4, 6.9, 8, 0.3, 9, RGBColor(0x64,0x74,0x8b), italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_BG)
rect(s, 0, 0, 0.12, 7.5, GREEN)
rect(s, 0.12, 0, 13.21, 0.8, GREEN)
txt(s, "What We'll Cover Today", 0.3, 0.12, 10, 0.6, 22, WHITE, bold=True)

items = [
    ("01", "The Problem",          "Why traditional club management fails",       "🔴"),
    ("02", "Our Solution",         "TrustNest — built for Herbalife clubs",        "🟢"),
    ("03", "Key Features",         "Member portal, barcode, invoices & more",      "🔵"),
    ("04", "Multi-Branch Power",   "Run multiple clubs from one platform",         "🟣"),
    ("05", "Analytics & Revenue",  "Data-driven decisions for growth",             "🟡"),
    ("06", "Getting Started",      "Setup in under 10 minutes",                    "⚪"),
]
for i,(num,title,sub,emoji) in enumerate(items):
    col = i % 2
    row = i // 2
    xl = 0.5 + col * 6.3
    yt = 1.05 + row * 1.85
    card(s, xl, yt, 5.9, 1.6, WHITE, SILVER)
    txt(s, num, xl+0.18, yt+0.12, 0.7, 0.5, 28, GREEN, bold=True)
    txt(s, emoji, xl+0.9, yt+0.12, 0.7, 0.5, 22, DARK)
    txt(s, title, xl+1.55, yt+0.12, 4.1, 0.45, 15, DARK, bold=True)
    txt(s, sub,   xl+1.55, yt+0.6,  4.1, 0.55, 10, SLATE, italic=True)

txt(s, "TrustNest HBL  •  Confidential", 0.3, 7.15, 8, 0.25, 8, SLATE, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 0.12, 7.5, RGBColor(0xEF,0x44,0x44))
rect(s, 0.12, 0, 13.21, 0.75, RGBColor(0x1e,0x1e,0x2e))
txt(s, "The Problem", 0.3, 0.1, 10, 0.55, 22, RGBColor(0xFC,0xA5,0xA5), bold=True)
txt(s, "Running a nutrition club today is harder than it should be", 0.3, 0.78, 12, 0.55, 18, WHITE, italic=True)

pains = [
    ("📋", "Paper Registers",    "Members sign in on paper — data gets lost, no history"),
    ("💸", "Manual Billing",     "Handwritten bills, no GST invoice, no PDF download"),
    ("📵", "No Digital Trail",   "Can't track who's active, who hasn't visited in days"),
    ("😓", "No Follow-Up System","Members drop off silently — no alerts, no WhatsApp nudge"),
    ("🏪", "Single-Club Chaos",  "Running 2+ clubs? Zero visibility across branches"),
    ("📊", "Blind Operations",   "No revenue reports, no low-stock alerts, no analytics"),
]
for i,(icon,title,desc) in enumerate(pains):
    col = i % 2; row = i // 2
    xl = 0.35 + col*6.45; yt = 1.45 + row*1.75
    rect(s, xl, yt, 6.1, 1.55, RGBColor(0x1e,0x14,0x14), RGBColor(0xEF,0x44,0x44), 0.4)
    txt(s, icon, xl+0.2, yt+0.3, 0.8, 0.7, 26, WHITE)
    txt(s, title, xl+1.0, yt+0.15, 4.8, 0.4, 14, RGBColor(0xFC,0xA5,0xA5), bold=True)
    txt(s, desc, xl+1.0, yt+0.62, 4.8, 0.6, 10, SILVER, italic=True)

txt(s, "TrustNest HBL  •  Confidential", 0.3, 7.15, 8, 0.25, 8, SLATE, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — OUR SOLUTION
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_BG)
rect(s, 0, 0, 0.12, 7.5, GREEN)
rect(s, 0.12, 0, 13.21, 0.75, GREEN)
txt(s, "The TrustNest Solution", 0.3, 0.1, 10, 0.55, 22, WHITE, bold=True)
txt(s, "One platform. Every club. Zero paperwork.", 0.3, 0.78, 12, 0.5, 17, DARK, italic=True)

# Central "hub" graphic simulation
rect(s, 4.0, 1.4, 5.33, 5.2, WHITE, GREEN, 0.8)
txt(s, "🌿", 5.9, 1.65, 1.5, 0.8, 40, GREEN, align=PP_ALIGN.CENTER)
txt(s, "TrustNest", 4.8, 2.5, 3.7, 0.6, 22, GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(s, "HBL Platform", 4.8, 3.05, 3.7, 0.45, 14, SLATE, align=PP_ALIGN.CENTER)

pillars = [
    (0.25, 1.4,  "📱 Member\nPortal"),
    (0.25, 3.1,  "🔍 Barcode\nCheck-In"),
    (0.25, 4.8,  "📄 Auto\nInvoice"),
    (9.95, 1.4,  "🏪 Multi\nBranch"),
    (9.95, 3.1,  "📊 Analytics\n& Reports"),
    (9.95, 4.8,  "💬 WhatsApp\nAlerts"),
]
for (xl, yt, label) in pillars:
    card(s, xl, yt, 2.2, 1.2, RGBColor(0xF0,0xFD,0xF4), GREEN)
    txt(s, label, xl+0.1, yt+0.2, 2.0, 0.9, 11, DARK, bold=True, align=PP_ALIGN.CENTER)

txt(s, "TrustNest HBL  •  Confidential", 0.3, 7.15, 8, 0.25, 8, SLATE, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MEMBER PORTAL FEATURES
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 0.12, 7.5, GREEN_MID)
rect(s, 0.12, 0, 13.21, 0.75, DARK)
txt(s, "Member Portal  —  Features at a Glance", 0.3, 0.1, 12, 0.55, 20, WHITE, bold=True)
txt(s, "Members download nothing. Just WhatsApp the link.", 0.3, 0.78, 12, 0.4, 13, SLATE, italic=True)

features = [
    ("🔐", "Phone OTP Login",        "One-tap login with 6-digit OTP. No password. No app download."),
    ("🪪", "Digital Membership Card","QR code + barcode card with plan tier, points, expiry date."),
    ("✅", "Daily Check-In",         "Members check in themselves — earns 1 loyalty point per visit."),
    ("🛒", "Shake & Supplement Order","Browse products, add to cart, choose pickup time, get PDF invoice."),
    ("⭐", "Loyalty Points",          "Auto-earn: 1pt check-in, 1pt per ₹100 spent. 10pts = free shake."),
    ("📜", "Order History",           "All past orders with PDF invoice download, inline view."),
    ("📸", "QR Auto-Login",           "Scan club QR → instantly logged in + checked in. Zero friction."),
    ("💳", "Plan Tiers",              "Basic / Silver / Gold / Platinum with colour-coded member cards."),
]
for i, (icon, title, desc) in enumerate(features):
    col = i % 2; row = i // 2
    xl = 0.35 + col * 6.45; yt = 1.28 + row * 1.45
    card(s, xl, yt, 6.1, 1.3, RGBColor(0xF8,0xFF,0xFA), SILVER)
    txt(s, icon, xl+0.18, yt+0.3, 0.7, 0.55, 22, DARK)
    txt(s, title, xl+1.0, yt+0.12, 4.8, 0.38, 13, DARK, bold=True)
    txt(s, desc,  xl+1.0, yt+0.55, 4.8, 0.55, 10, SLATE, italic=True)

txt(s, "TrustNest HBL  •  Confidential", 0.3, 7.15, 8, 0.25, 8, SLATE, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — BARCODE & QR CHECK-IN
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 0.12, 7.5, TEAL)
rect(s, 0.12, 0, 13.21, 0.75, RGBColor(0x0a,0x1a,0x2a))
txt(s, "Barcode & QR Check-In System", 0.3, 0.1, 12, 0.55, 20, WHITE, bold=True)
txt(s, "Admin scans, system records — no pen, no paper, no mistakes.", 0.3, 0.78, 12, 0.4, 13, RGBColor(0x99,0xe6,0xca), italic=True)

# Two column layout
rect(s, 0.35, 1.25, 5.9, 5.5, RGBColor(0x0d,0x23,0x2e), RGBColor(0x0D,0x94,0x88), 0.5)
txt(s, "👤  Member Card", 0.55, 1.38, 5.5, 0.45, 14, TEAL, bold=True)
txt(s, "Every member gets a unique card with:", 0.55, 1.88, 5.5, 0.35, 11, SILVER)
items_l = [
    "QR Code — deep-link login + auto check-in",
    "CODE-128 Barcode — for admin scanner",
    "Member code  e.g.  VPT-0001",
    "Plan badge, points, expiry date",
    "Save / Share / Print actions",
]
for i,t in enumerate(items_l):
    txt(s, f"  ▶  {t}", 0.55, 2.3+i*0.52, 5.5, 0.45, 10.5, WHITE)

rect(s, 6.75, 1.25, 6.2, 5.5, RGBColor(0x0d,0x23,0x2e), RGBColor(0x0D,0x94,0x88), 0.5)
txt(s, "🔍  Admin Scanner", 6.95, 1.38, 5.8, 0.45, 14, TEAL, bold=True)
txt(s, "Admin opens Scanner on any phone/tablet:", 6.95, 1.88, 5.8, 0.35, 11, SILVER)
steps = [
    ("1", "Tap 'Start Scanner' — camera opens"),
    ("2", "Point at member's barcode or QR code"),
    ("3", "System identifies member instantly"),
    ("4", "Check-in recorded + 1 loyalty point awarded"),
    ("5", "Member details shown on screen"),
    ("     ", "OR type member code manually as fallback"),
]
for i,(num,step) in enumerate(steps):
    clr = TEAL if num.strip().isdigit() else SILVER
    txt(s, f"  {num}.  {step}", 6.95, 2.3+i*0.52, 5.8, 0.45, 10.5, clr if num.strip().isdigit() else RGBColor(0x9c,0xa3,0xaf))

txt(s, "TrustNest HBL  •  Confidential", 0.3, 7.15, 8, 0.25, 8, SLATE, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TAX INVOICE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_BG)
rect(s, 0, 0, 0.12, 7.5, GREEN)
rect(s, 0.12, 0, 13.21, 0.75, GREEN)
txt(s, "Auto-Generated GST Tax Invoice", 0.3, 0.1, 12, 0.55, 20, WHITE, bold=True)
txt(s, "Professional PDF invoice — auto downloads right after checkout.", 0.3, 0.78, 12, 0.4, 13, DARK, italic=True)

# Invoice mock
rect(s, 0.35, 1.2, 6.0, 5.6, WHITE, GREEN, 0.8)
rect(s, 0.35, 1.2, 6.0, 0.75, GREEN)
txt(s, "Herbalife  TAX INVOICE", 0.5, 1.3, 5.6, 0.5, 13, WHITE, bold=True)
txt(s, "Original for Recipient", 0.5, 1.72, 5.6, 0.25, 8, RGBColor(0xbb,0xff,0xcc))
rect(s, 0.35, 1.95, 6.0, 0.6, RGBColor(0xf0,0xfd,0xf4))
txt(s, "Club Name  |  GSTIN  |  FSSAI No.", 0.5, 2.02, 5.6, 0.3, 8.5, SLATE)
txt(s, "Invoice No:  INV-HBL-20260601-4231", 0.5, 2.55, 5.6, 0.28, 8, DARK)
txt(s, "Member: Ravi Kumar  |  Ph: 9876543210", 0.5, 2.85, 5.6, 0.28, 8, DARK)
# Mini table header
rect(s, 0.35, 3.18, 6.0, 0.35, GREEN)
for col,lbl in zip([0,1.5,3.5,4.5],[" Item","MRP","Retail","GST"]):
    txt(s, lbl, 0.45+col, 3.22, 1.4, 0.28, 7.5, WHITE, bold=True)
for i,(item,mrp,ret,gst) in enumerate([
    ("Formula 1 Shake","315","283","14"),("Afresh Lemon Tea","165","148","7"),("Cell Activator","221","210","10"),
]):
    yt = 3.56+i*0.38
    bg_r = RGBColor(0xf8,0xff,0xfa) if i%2==0 else WHITE
    rect(s, 0.35, yt, 6.0, 0.37, bg_r)
    for col,val in zip([0,1.5,3.5,4.5],[item,mrp,ret,gst]):
        txt(s, val, 0.45+col, yt+0.05, 1.4, 0.28, 7.5, DARK)
rect(s, 0.35, 4.7, 6.0, 0.45, RGBColor(0xf0,0xfd,0xf4))
txt(s, "Invoice Total: Rs. 641.00", 3.0, 4.76, 3.2, 0.33, 11, GREEN, bold=True)
txt(s, "Net Savings vs MRP: Rs. 60.00", 0.45, 4.76, 3.2, 0.33, 9, SLATE)
txt(s, "Six Hundred Forty One Rupees Only", 0.45, 5.22, 5.5, 0.3, 7.5, SLATE, italic=True)

# Benefits list on right
rect(s, 6.75, 1.2, 6.2, 5.6, WHITE, SILVER, 0.5)
txt(s, "📄  What's Included in Every Invoice", 6.9, 1.3, 5.9, 0.45, 13, DARK, bold=True)
points = [
    ("✅", "Herbalife branding header (Green)"),
    ("✅", "Club name, GSTIN, FSSAI number"),
    ("✅", "Member details + shipping info"),
    ("✅", "SKU, MRP, Retail price, Discount"),
    ("✅", "Taxable value + SGST 2.5% + CGST 2.5%"),
    ("✅", "Invoice total in words (GST compliant)"),
    ("✅", "Volume Points for distributor tracking"),
    ("✅", "Authorised digital signatory block"),
    ("✅", "Auto-download on checkout"),
    ("✅", "Inline PDF view option"),
]
for i,(ic,pt) in enumerate(points):
    txt(s, f"{ic}  {pt}", 6.9, 1.85+i*0.47, 5.9, 0.4, 11, DARK if ic=="✅" else SLATE)

txt(s, "TrustNest HBL  •  Confidential", 0.3, 7.15, 8, 0.25, 8, SLATE, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ADMIN DASHBOARD
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 0.12, 7.5, GREEN_MID)
rect(s, 0.12, 0, 13.21, 0.75, RGBColor(0x10,0x20,0x18))
txt(s, "Admin Dashboard — Full Control", 0.3, 0.1, 12, 0.55, 20, WHITE, bold=True)
txt(s, "Everything you need to run the club, in one screen.", 0.3, 0.78, 12, 0.4, 13, GREEN_MID, italic=True)

kpis = [
    ("48", "Active\nMembers", GREEN_MID),
    ("12", "Today's\nCheck-Ins", TEAL),
    ("7",  "Orders\nToday",    AMBER),
    ("3",  "Renewals\nDue",    RGBColor(0xF8,0x71,0x71)),
]
for i,(val,lbl,clr) in enumerate(kpis):
    xl = 0.35 + i*3.05
    rect(s, xl, 1.25, 2.75, 1.45, RGBColor(0x12,0x24,0x18), clr, 0.5)
    txt(s, val, xl+0.2, 1.32, 2.4, 0.75, 38, clr, bold=True, align=PP_ALIGN.CENTER)
    txt(s, lbl, xl+0.2, 2.0, 2.4, 0.55, 10, SILVER, align=PP_ALIGN.CENTER)

modules = [
    ("👥", "Members",       "Add, search, filter, renew memberships"),
    ("🔍", "Scanner",       "Barcode / QR check-in from any phone"),
    ("🆕", "New Joiners",   "Day counter, milestones, 30-day journey"),
    ("🥤", "Shake Tracker", "7-day activity grid, active vs. at-risk"),
    ("📞", "Follow-Ups",    "HIGH/MED/LOW priority WhatsApp outreach"),
    ("📦", "Products",      "Stock levels, low-stock alerts, pricing"),
    ("🧾", "Orders",        "Live kitchen queue, status management"),
    ("📊", "Reports",       "Revenue, popular products, retention charts"),
    ("⚙️", "Settings",      "Branch name, GSTIN, FSSAI, PIN change"),
]
for i,(icon,name,desc) in enumerate(modules):
    col = i % 3; row = i // 3
    xl = 0.35 + col*4.25; yt = 2.9 + row*1.4
    rect(s, xl, yt, 4.0, 1.25, RGBColor(0x0e,0x1e,0x14), GREEN, 0.3)
    txt(s, icon, xl+0.18, yt+0.2, 0.7, 0.65, 18, WHITE)
    txt(s, name, xl+0.95, yt+0.15, 2.8, 0.38, 12, GREEN_MID, bold=True)
    txt(s, desc, xl+0.95, yt+0.58, 2.8, 0.45, 9, SILVER, italic=True)

txt(s, "TrustNest HBL  •  Confidential", 0.3, 7.15, 8, 0.25, 8, SLATE, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MULTI-TENANT / MULTI-BRANCH
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_BG)
rect(s, 0, 0, 0.12, 7.5, RGBColor(0x76,0x29,0xD8))
rect(s, 0.12, 0, 13.21, 0.75, RGBColor(0x2d,0x1b,0x69))
txt(s, "Multi-Branch Architecture", 0.3, 0.1, 12, 0.55, 20, WHITE, bold=True)
txt(s, "Scale from 1 club to 50 clubs — same login, separate data.", 0.3, 0.78, 12, 0.4, 13, RGBColor(0x76,0x29,0xD8), italic=True)

# Hub diagram
rect(s, 4.92, 2.9, 3.5, 1.5, RGBColor(0x2d,0x1b,0x69), RGBColor(0x76,0x29,0xD8), 1.0)
txt(s, "🌿 TrustNest", 4.92, 3.0, 3.5, 0.5, 15, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "Central Platform", 4.92, 3.45, 3.5, 0.4, 10, RGBColor(0xc4,0xb5,0xfd), align=PP_ALIGN.CENTER)

branches = [
    (0.5,  1.1, "🏪 Veppampattu\nNutrition Club",    "PIN: 1234"),
    (0.5,  4.2, "🏪 Chennai Central\nNutrition Club", "PIN: 5678"),
    (9.3,  1.1, "🏪 Tambaram\nNutrition Club",        "PIN: 9012"),
    (9.3,  4.2, "🏪 Coimbatore\nNutrition Club",      "PIN: 3456"),
]
for (xl,yt,name,pin) in branches:
    card(s, xl, yt, 3.5, 1.7, WHITE, RGBColor(0x76,0x29,0xD8))
    txt(s, name, xl+0.15, yt+0.18, 3.2, 0.65, 11, DARK, bold=True)
    txt(s, pin, xl+0.15, yt+0.92, 3.2, 0.35, 10, SLATE)
    txt(s, "Isolated data ✓", xl+0.15, yt+1.28, 3.2, 0.3, 9, GREEN, italic=True)

benefits = [
    "Each branch has its own admin PIN",
    "Members, products, orders 100% isolated",
    "Custom branch name on every invoice",
    "Admin sees only their club's data",
    "Super-admin can view all branches",
    "Branch URL: /hbl?t=slug shared with members",
]
rect(s, 0.35, 6.1, 12.6, 1.1, WHITE, SILVER, 0.3)
txt(s, "Isolation Guarantees:", 0.55, 6.18, 2.5, 0.35, 11, DARK, bold=True)
for i,b in enumerate(benefits):
    col = i % 3; row = i // 3
    txt(s, f"✦  {b}", 3.1+col*3.25, 6.18+row*0.42, 3.1, 0.38, 9.5, SLATE)

txt(s, "TrustNest HBL  •  Confidential", 0.3, 7.15, 8, 0.25, 8, SLATE, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — FOLLOW-UP & RETENTION ENGINE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 0.12, 7.5, AMBER)
rect(s, 0.12, 0, 13.21, 0.75, DARK)
txt(s, "Retention Engine — Never Lose a Member", 0.3, 0.1, 12, 0.55, 20, WHITE, bold=True)
txt(s, "Automated follow-up intelligence powered by real visit data.", 0.3, 0.78, 12, 0.4, 13, AMBER, italic=True)

cols_data = [
    (RGBColor(0xFF,0xE4,0xE6), RGBColor(0xEF,0x44,0x44), "HIGH Priority 🔴",
     ["No check-in for 7+ days","Membership expires in ≤2 days","Member never visited"],
     "WhatsApp Template:\n\"Hi [Name]! We miss you at the club.\nYour membership expires soon — let's\nrenew and get back on track! 💚\""),
    (RGBColor(0xFE,0xF3,0xC7), AMBER, "MEDIUM Priority 🟡",
     ["No check-in 3–7 days","Expiry in 3–7 days","Day 3 / 7 milestone"],
     "WhatsApp Template:\n\"Hi [Name]! Your consistency is key.\nCome in today for your nutrition shake\nand earn loyalty points! ⭐\""),
    (RGBColor(0xF0,0xFD,0xF4), GREEN, "LOW Priority 🟢",
     ["No order in 14 days","New joiner day 14 / 30","Occasional visitor"],
     "WhatsApp Template:\n\"Hi [Name]! Did you know our Formula 1\nShake helps with your goal? Order now\nand we'll have it ready! 🥤\""),
]
for i,(bg_c,ac,title,pts,wa) in enumerate(cols_data):
    xl = 0.35 + i*4.25
    rect(s, xl, 1.25, 4.0, 5.55, bg_c, ac, 0.5)
    txt(s, title, xl+0.2, 1.35, 3.6, 0.45, 12, DARK, bold=True)
    txt(s, "Triggers:", xl+0.2, 1.88, 3.6, 0.3, 10, DARK, bold=True)
    for j,pt in enumerate(pts):
        txt(s, f"  • {pt}", xl+0.2, 2.2+j*0.42, 3.6, 0.38, 10, DARK)
    txt(s, wa, xl+0.2, 3.55, 3.6, 1.5, 8.5, SLATE, italic=True)
    txt(s, "📲 One-tap WA send", xl+0.2, 6.35, 3.6, 0.32, 10, ac, bold=True)

txt(s, "TrustNest HBL  •  Confidential", 0.3, 7.15, 8, 0.25, 8, SLATE, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — ANALYTICS & REPORTS
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 0.12, 7.5, TEAL)
rect(s, 0.12, 0, 13.21, 0.75, RGBColor(0x0a,0x1a,0x2a))
txt(s, "Analytics & Revenue Reports", 0.3, 0.1, 12, 0.55, 20, WHITE, bold=True)
txt(s, "Know your numbers. Grow with confidence.", 0.3, 0.78, 12, 0.4, 13, TEAL, italic=True)

# Bar chart simulation (daily revenue)
chart_data = [1200,1800,2400,1600,3100,2700,3800,2200,4100,3500,2900,4800,3200,4500]
chart_max = max(chart_data)
chart_l, chart_t, chart_w, chart_h = 0.4, 1.25, 6.0, 4.2
rect(s, chart_l, chart_t, chart_w, chart_h, RGBColor(0x0d,0x1e,0x2a), TEAL, 0.3)
txt(s, "Daily Revenue (Rs.) — This Month", chart_l, chart_t+0.08, chart_w, 0.3, 10, TEAL, bold=True)
bar_w = 0.35
for i,val in enumerate(chart_data):
    bh = (val/chart_max) * 3.2
    bx = chart_l + 0.2 + i*(bar_w+0.08)
    by = chart_t + chart_h - 0.25 - bh
    clr = GREEN_MID if val == max(chart_data) else TEAL
    rect(s, bx, by, bar_w, bh, clr)

# Right panel — metrics
rect(s, 6.8, 1.25, 6.15, 4.2, RGBColor(0x0d,0x1e,0x2a), TEAL, 0.3)
txt(s, "📊  Key Metrics This Month", 7.0, 1.35, 5.8, 0.38, 12, TEAL, bold=True)
metrics = [
    ("Rs. 48,600", "Total Revenue", GREEN_MID),
    ("312",        "Total Orders",  WHITE),
    ("89%",        "Retention Rate", GREEN_MID),
    ("23",         "New Members",   TEAL),
    ("Formula 1", "Top Product",   WHITE),
    ("Gold",       "Top Plan Tier", AMBER),
]
for i,(val,lbl,clr) in enumerate(metrics):
    row=i//2; col=i%2
    xt = 7.0 + col*3.0; yt = 1.85 + row*1.1
    txt(s, val, xt, yt, 2.7, 0.5, 18, clr, bold=True)
    txt(s, lbl, xt, yt+0.45, 2.7, 0.35, 9.5, SILVER)

# Bottom strip
rect(s, 0.4, 5.7, 12.55, 1.1, RGBColor(0x0d,0x20,0x18), GREEN, 0.3)
reports = ["Week / Month / Year range selector","Daily revenue bar chart","Popular products ranking","Member retention donut chart","New joiner growth trend"]
txt(s, "Report Views Available:", 0.6, 5.77, 2.5, 0.38, 10, GREEN_MID, bold=True)
for i,r in enumerate(reports):
    txt(s, f"  ✦  {r}", 3.1+i*1.9, 5.77, 1.85, 0.75, 8.5, SILVER)

txt(s, "TrustNest HBL  •  Confidential", 0.3, 7.15, 8, 0.25, 8, SLATE, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SHAKE TRACKER + NEW JOINERS
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_BG)
rect(s, 0, 0, 0.12, 7.5, GREEN)
rect(s, 0.12, 0, 13.21, 0.75, GREEN)
txt(s, "Shake Tracker & New Joiner Journey", 0.3, 0.1, 12, 0.55, 20, WHITE, bold=True)
txt(s, "Track every member's progress — from day 1 to day 30.", 0.3, 0.78, 12, 0.4, 13, DARK, italic=True)

# Left: Shake Tracker
rect(s, 0.35, 1.2, 5.9, 5.6, WHITE, SILVER, 0.4)
txt(s, "🥤  Shake Tracker  —  7-Day Grid", 0.55, 1.3, 5.6, 0.4, 13, DARK, bold=True)
headers = ["Mon","Tue","Wed","Thu","Fri","Sat","Sun"]
members_t = [("Ravi Kumar","✅","✅","✅","✅","✅","⬜","✅"),
             ("Priya S.",  "✅","⬜","✅","✅","⬜","⬜","✅"),
             ("Karthik M.","✅","✅","⬜","⬜","⬜","⬜","⬜"),
             ("Meena R.",  "✅","✅","✅","✅","✅","✅","✅")]
for j,h in enumerate(headers):
    txt(s, h, 0.55+0.7+j*0.68, 1.8, 0.65, 0.28, 8.5, SLATE, bold=True, align=PP_ALIGN.CENTER)
for i,(name,*days) in enumerate(members_t):
    yt = 2.15+i*0.75
    clr = GREEN if all(d=="✅" for d in days) else (RGBColor(0xFF,0x80,0x80) if days.count("✅")<=2 else DARK)
    txt(s, name, 0.55, yt+0.05, 1.5, 0.38, 9.5, clr, bold=True)
    for j,d in enumerate(days):
        txt(s, d, 0.55+0.7+j*0.68, yt, 0.65, 0.42, 13, DARK, align=PP_ALIGN.CENTER)

# Right: New Joiners
rect(s, 6.75, 1.2, 6.2, 5.6, WHITE, SILVER, 0.4)
txt(s, "🆕  30-Day Journey Milestones", 6.95, 1.3, 5.9, 0.4, 13, DARK, bold=True)
milestones = [
    ("Day 1",   "Joined",             "🟢", "Welcome WhatsApp sent"),
    ("Day 3",   "First shake check",  "🟡", "Engagement follow-up"),
    ("Day 7",   "One week in",        "🏆", "7-day streak badge"),
    ("Day 14",  "2-week milestone",   "⭐", "Progress check-in"),
    ("Day 21",  "Habit forming",      "🔥", "Consistency badge"),
    ("Day 30",  "30-day champion",    "💎", "Platinum loyalty gift"),
]
for i,(day,event,badge,action) in enumerate(milestones):
    yt = 1.85+i*0.77
    rect(s, 6.95, yt, 5.8, 0.65, RGBColor(0xf8,0xff,0xfa), GREEN, 0.3)
    txt(s, badge, 7.08, yt+0.1, 0.55, 0.42, 18, DARK)
    txt(s, f"{day}: {event}", 7.7, yt+0.08, 3.0, 0.3, 10.5, DARK, bold=True)
    txt(s, action, 7.7, yt+0.38, 3.0, 0.25, 9, SLATE, italic=True)

txt(s, "TrustNest HBL  •  Confidential", 0.3, 7.15, 8, 0.25, 8, SLATE, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — HOW IT WORKS (FLOW)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 0.12, 7.5, GREEN_MID)
rect(s, 0.12, 0, 13.21, 0.75, RGBColor(0x0a,0x18,0x10))
txt(s, "How It Works — From Setup to Daily Operations", 0.3, 0.1, 12, 0.55, 20, WHITE, bold=True)
txt(s, "Your club goes digital in under 10 minutes.", 0.3, 0.78, 12, 0.4, 13, GREEN_MID, italic=True)

steps = [
    ("1", "Admin Setup\n(10 min)", "POST /api/setup with club name,\nGSTIN, FSSAI, admin PIN", "🛠️"),
    ("2", "Add Members\n(30 sec each)", "Name + Phone → system generates\nQR + barcode member card", "👤"),
    ("3", "Load Products\n(1 click)", "12 default Herbalife products\nwith SKU, MRP, pricing — instant", "📦"),
    ("4", "Daily Check-In\n(2 sec)", "Admin scans barcode OR\nmember self-checks via QR link", "✅"),
    ("5", "Member Orders\n(1 min)", "Browse → cart → checkout →\nPDF invoice auto-downloads", "🛒"),
    ("6", "Admin Reviews\n(5 min/day)", "Dashboard KPIs, follow-up list,\nWhatsApp alerts with 1 tap", "📊"),
]
arrow_y = 3.25
for i,(num,title,desc,ico) in enumerate(steps):
    xl = 0.25 + i*2.18
    rect(s, xl, 1.25, 1.95, 2.5, RGBColor(0x0e,0x22,0x18), GREEN, 0.5)
    txt(s, ico,   xl+0.7, 1.38, 0.8, 0.6, 22, WHITE, align=PP_ALIGN.CENTER)
    txt(s, num,   xl+0.1, 1.98, 1.75, 0.5, 28, GREEN_MID, bold=True, align=PP_ALIGN.CENTER)
    txt(s, title, xl+0.05, 2.45, 1.85, 0.55, 10, WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, desc,  xl+0.05, 3.05, 1.85, 0.65, 8, SILVER, italic=True, align=PP_ALIGN.CENTER)
    if i < 5:
        txt(s, "→", xl+2.0, arrow_y, 0.2, 0.4, 18, GREEN_MID, bold=True)

# Bottom: For member journey
rect(s, 0.35, 4.0, 12.6, 2.95, RGBColor(0x0a,0x18,0x10), GREEN, 0.3)
txt(s, "Member Journey:", 0.6, 4.12, 3.0, 0.38, 11, GREEN_MID, bold=True)
journey = [
    ("📱", "Gets club URL", "/hbl?t=veppampattu"),
    ("🔢", "Enters phone", "OTP sent instantly"),
    ("🏠", "Sees dashboard", "Card, points, QR"),
    ("🛒", "Places order",  "PDF invoice ready"),
    ("📸", "Shows QR",      "Admin scans, checked in"),
    ("💬", "Gets WA nudge", "On milestones & expiry"),
]
for i,(ico,step,sub) in enumerate(journey):
    xl = 0.6 + i*2.08
    txt(s, ico, xl, 4.55, 0.6, 0.5, 22, WHITE)
    txt(s, step, xl+0.55, 4.55, 1.5, 0.35, 10, WHITE, bold=True)
    txt(s, sub,  xl+0.55, 4.92, 1.5, 0.35, 8.5, SILVER, italic=True)
    if i < 5:
        txt(s, "→", xl+2.0, 4.65, 0.2, 0.38, 14, GREEN_MID, bold=True)

txt(s, "TrustNest HBL  •  Confidential", 0.3, 7.15, 8, 0.25, 8, SLATE, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — WHY TRUSTNEST (COMPETITIVE)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, LIGHT_BG)
rect(s, 0, 0, 0.12, 7.5, GREEN)
rect(s, 0.12, 0, 13.21, 0.75, GREEN)
txt(s, "Why TrustNest?  —  vs. The Alternatives", 0.3, 0.1, 12, 0.55, 20, WHITE, bold=True)
txt(s, "Purpose-built for Herbalife clubs. Nothing else comes close.", 0.3, 0.78, 12, 0.4, 13, DARK, italic=True)

# Comparison table
headers_row = ["Feature", "Paper / WhatsApp", "Generic POS", "🌿 TrustNest"]
col_widths = [3.2, 2.5, 2.5, 3.0]
col_starts = [0.35, 3.55, 6.05, 8.55]
colors_h   = [DARK, RGBColor(0x64,0x74,0x8b), RGBColor(0x64,0x74,0x8b), GREEN]
bg_h       = [SILVER, SILVER, SILVER, GREEN]

rect(s, 0.35, 1.25, 11.2, 0.55, GREEN)
for ci,(h,xl,wid) in enumerate(zip(headers_row,col_starts,col_widths)):
    txt(s, h, xl+0.1, 1.33, wid-0.2, 0.38, 11, WHITE, bold=True)

rows = [
    ("GST Tax Invoice (PDF)",           "❌","❌","✅"),
    ("QR + Barcode Check-In",           "❌","❌","✅"),
    ("Multi-Branch Support",            "❌","💰 Extra","✅"),
    ("WhatsApp Follow-Up Automation",   "Manual","❌","✅"),
    ("Loyalty Points + Rewards",        "❌","❌","✅"),
    ("Herbalife SKU Catalog Built-In",  "❌","❌","✅"),
    ("Shake / Activity Tracker",        "❌","❌","✅"),
    ("New Joiner Journey (30 days)",    "❌","❌","✅"),
    ("Mobile-First, No App Install",    "❌","❌","✅"),
    ("Monthly Cost",                    "₹0", "₹2,000–5,000","Ask us"),
]
for ri,row in enumerate(rows):
    yt = 1.82 + ri*0.48
    bg_c = WHITE if ri%2==0 else RGBColor(0xf8,0xff,0xfa)
    rect(s, 0.35, yt, 11.2, 0.46, bg_c, SILVER, 0.2)
    clrs = [DARK, RGBColor(0xEF,0x44,0x44), RGBColor(0xEF,0x44,0x44), GREEN]
    for ci,(val,xl,wid) in enumerate(zip(row,col_starts,col_widths)):
        clr = clrs[ci] if ci > 0 else DARK
        bld = ci == 0 or ci == 3
        txt(s, val, xl+0.1, yt+0.08, wid-0.2, 0.33, 10.5, clr, bold=bld)

txt(s, "TrustNest HBL  •  Confidential", 0.3, 7.15, 8, 0.25, 8, SLATE, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — GETTING STARTED
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 0.12, 7.5, GREEN_MID)
rect(s, 0.12, 0, 13.21, 0.75, RGBColor(0x08,0x18,0x10))
txt(s, "Getting Started", 0.3, 0.1, 12, 0.55, 20, WHITE, bold=True)
txt(s, "Your club can be live today — here's exactly what to do.", 0.3, 0.78, 12, 0.4, 13, GREEN_MID, italic=True)

steps_gs = [
    ("Step 1", "Share your club details",
     "Name, Address, GSTIN, FSSAI, Phone number\nAdmin PIN of your choice",
     "📋", RGBColor(0x0e,0x28,0x18)),
    ("Step 2", "We create your branch in 5 min",
     "Unique branch URL:  trustnest.com/hbl?t=yourclub\nAll 12 default products pre-loaded",
     "⚡", RGBColor(0x0e,0x24,0x1a)),
    ("Step 3", "Add your members",
     "Name + phone per member\nBarcode + QR card auto-generated for each",
     "👥", RGBColor(0x0a,0x20,0x14)),
    ("Step 4", "Share the link with members",
     "WhatsApp the URL — members log in instantly\nNo app install, no training needed",
     "💬", RGBColor(0x08,0x1e,0x12)),
]
for i,(step,title,desc,ico,bg_c) in enumerate(steps_gs):
    xl = 0.35 + (i%2)*6.45; yt = 1.25 + (i//2)*2.6
    rect(s, xl, yt, 6.1, 2.38, bg_c, GREEN, 0.5)
    txt(s, ico,   xl+0.2, yt+0.25, 0.8, 0.7, 30, WHITE)
    txt(s, step,  xl+1.1, yt+0.15, 4.7, 0.35, 10, GREEN_MID, bold=True)
    txt(s, title, xl+1.1, yt+0.55, 4.7, 0.45, 15, WHITE, bold=True)
    txt(s, desc,  xl+1.1, yt+1.1,  4.7, 0.9,  10, SILVER, italic=True)

txt(s, "TrustNest HBL  •  Confidential", 0.3, 7.15, 8, 0.25, 8, SLATE, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CLOSING / CALL TO ACTION
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK)
rect(s, 0, 0, 13.33, 7.5, GREEN)
# Dark overlay
rect(s, 0, 0, 13.33, 7.5, DARK)

# Large decorative circles
c = s.shapes.add_shape(9, Inches(-1), Inches(-1), Inches(6), Inches(6))
c.fill.solid(); c.fill.fore_color.rgb = GREEN; c.line.fill.background()
sp = c._element; spPr = sp.find(nsmap.qn('p:spPr'))
sf = spPr.find('.//' + nsmap.qn('a:solidFill'))
if sf:
    sr = sf.find(nsmap.qn('a:srgbClr'))
    if sr is not None:
        a = etree.SubElement(sr, nsmap.qn('a:alpha')); a.set('val','8000')

c2 = s.shapes.add_shape(9, Inches(9), Inches(3.5), Inches(5), Inches(5))
c2.fill.solid(); c2.fill.fore_color.rgb = TEAL; c2.line.fill.background()
sp2 = c2._element; spPr2 = sp2.find(nsmap.qn('p:spPr'))
sf2 = spPr2.find('.//' + nsmap.qn('a:solidFill'))
if sf2:
    sr2 = sf2.find(nsmap.qn('a:srgbClr'))
    if sr2 is not None:
        a2 = etree.SubElement(sr2, nsmap.qn('a:alpha')); a2.set('val','8000')

txt(s, "🌿", 5.9, 0.7, 1.5, 1.0, 44, GREEN_MID, bold=True, align=PP_ALIGN.CENTER)
txt(s, "Ready to Transform Your Club?", 1.0, 1.7, 11.33, 1.1, 36, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "TrustNest HBL is live. Your branch can be set up today.", 1.0, 2.85, 11.33, 0.55, 18, GREEN_MID, italic=True, align=PP_ALIGN.CENTER)

# 3 CTA boxes
cta = [
    ("🌐", "Visit the Platform", "trustnest-tsgz.vercel.app/hbl"),
    ("📲", "WhatsApp Us",        "wa.me/91XXXXXXXXXX"),
    ("📧", "Email",              "hello@trustnest.in"),
]
for i,(ico,lbl,val) in enumerate(cta):
    xl = 1.3 + i*3.6
    rect(s, xl, 3.65, 3.2, 1.65, RGBColor(0x0e,0x22,0x18), GREEN_MID, 0.5)
    txt(s, ico, xl+0.25, 3.78, 0.7, 0.55, 24, WHITE)
    txt(s, lbl, xl+1.05, 3.82, 2.0, 0.38, 11.5, WHITE, bold=True)
    txt(s, val, xl+1.05, 4.25, 2.0, 0.35, 9,    GREEN_MID, italic=True)

rect(s, 2.5, 5.55, 8.33, 0.85, GREEN, None, 0)
txt(s, "\"The club that tracks, grows.  The club that ignores, loses.\"", 1.0, 5.58, 11.33, 0.7, 14, WHITE, italic=True, align=PP_ALIGN.CENTER)

txt(s, "TrustNest  •  Built for India's Herbalife Clubs  •  2026", 1.0, 6.65, 11.33, 0.35, 9, RGBColor(0x9c,0xa3,0xaf), align=PP_ALIGN.CENTER)

# ─── Save ──────────────────────────────────────────────────────────────────
out = "/home/user/trustnest/TrustNest_HBL_Pitch_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
